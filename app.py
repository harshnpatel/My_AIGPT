import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter 
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from dotenv import load_dotenv
import requests
import json
import time
import datetime
import uuid
import pandas as pd
from pptx import Presentation
import shutil # Import shutil for directory removal

load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# --- Constants ---
FAISS_INDICES_DIR = "faiss_indices" # Directory to store individual FAISS indices

# --- Helper Functions for Document Processing ---
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        # Ensure the file pointer is at the beginning
        pdf.seek(0) 
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks, session_id):
    """
    Creates and saves a FAISS vector store unique to the given session_id.
    Returns the path to the saved index.
    """
    if not os.path.exists(FAISS_INDICES_DIR):
        os.makedirs(FAISS_INDICES_DIR)
    
    index_path = os.path.join(FAISS_INDICES_DIR, f"{session_id}_index")
    
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local(index_path)
    return index_path # Return the path where it was saved

# --- LLM Response Generation (unchanged) ---
def get_llm_response(user_question, context_text, selected_model, vector_store_available):
    if selected_model == "Llama 3.2 (Ollama)":
        return get_ollama_llama_response(user_question, context_text, vector_store_available)
    elif selected_model == "Gemini API":
        return get_gemini_response(user_question, context_text, vector_store_available)
    else:
        raise ValueError(f"Invalid model selection: {selected_model}")

def get_ollama_llama_response(user_question, context_text, vector_store_available):
    ollama_url = "http://localhost:11434/api/generate"
    ollama_model = "llama3.2"

    if vector_store_available and context_text.strip():
        prompt = f"""
        You are a highly knowledgeable AI assistant.
        Your primary goal is to answer the user's question accurately.

        Here is some potentially relevant document context:
        ---
        {context_text}
        ---

        Instructions:
        1.  **Prioritize the provided context:** If the answer is clearly and fully present in the context, answer *only* from the context.
        2.  **State if not in context:** If the answer is *not* available or only partially available in the provided context, clearly state "The answer is not available in the provided context." and then proceed to answer using your general knowledge if possible.
        3.  **Do not invent information:** Do not provide details that are not present in the context or your general knowledge.
        4.  **For general questions:** If the question is clearly a general knowledge question (e.g., "What is the capital of France?", "Tell me a joke") and the provided context is irrelevant, ignore the context and answer from your general knowledge.

        Question:
        {user_question}

        Answer:
        """
    else:
        prompt = f"""
        You are a helpful AI assistant. Answer the following question directly and comprehensively using your general knowledge.

        Question:
        {user_question}

        Answer:
        """

    payload = {
        "model": ollama_model,
        "prompt": prompt,
        "stream": True,
    }

    try:
        response = requests.post(ollama_url, json=payload, stream=True)
        response.raise_for_status()

        def generate_response_content():
            for chunk in response.iter_lines():
                if chunk:
                    try:
                        json_chunk = json.loads(chunk)
                        if "response" in json_chunk:
                            yield json_chunk["response"]
                        elif "message" in json_chunk and "content" in json_chunk["message"]:
                            yield json_chunk["message"]["content"]
                    except json.JSONDecodeError:
                        yield "Error decoding chunk."

        return generate_response_content()

    except requests.exceptions.ConnectionError:
        st.error(
            f"Could not connect to Ollama. Please ensure Ollama is running and '{ollama_model}' is pulled."
        )
        return ["Sorry, I can't connect to the local AI model (Ollama). Is it running?"] # Return as a list for consistent iteration
    except requests.exceptions.RequestException as e:
        st.error(f"Error querying Ollama: {e}. Check if model '{ollama_model}' is available.")
        print(f"Ollama Request Error: {e}, Response: {response.text if 'response' in locals() else 'N/A'}")
        return ["Sorry, there was an error with the local AI model."] # Return as a list

def get_gemini_response(user_question, context_text, vector_store_available):
    model_name = "gemini-1.5-flash"

    try:
        model = genai.GenerativeModel(model_name)
    except Exception as e:
        st.error(
            f"Could not load model '{model_name}'. Check your API key and model availability. Error: {e}")
        return "Sorry, I couldn't connect to the AI model."

    if vector_store_available and context_text.strip():
        prompt = f"""
        You are a highly knowledgeable AI assistant.
        Your primary goal is to answer the user's question accurately.

        Here is some potentially relevant document context:
        ---
        {context_text}
        ---

        Instructions:
        1.  **Prioritize the provided context:** If the answer is clearly and fully present in the context, answer *only* from the context.
        2.  **State if not in context:** If the answer is *not* available or only partially available in the provided context, clearly state "The answer is not available in the provided context." and then proceed to answer using your general knowledge if possible.
        3.  **Do not invent information:** Do not provide details that are not present in the context or your general knowledge.
        4.  **For general questions:** If the question is clearly a general knowledge question (e.g., "What is the capital of France?", "Tell me a joke") and the provided context is irrelevant, ignore the context and answer from your general knowledge.

        Question:
        {user_question}

        Answer:
        """
    else:
        prompt = f"""
        You are a helpful AI assistant. Answer the following question directly and comprehensively using your general knowledge.

        Question:
        {user_question}

        Answer:
        """

    try:
        response = model.generate_content(prompt, stream=True) # Enable streaming for Gemini
        return (chunk.text for chunk in response) # Return a generator for Gemini
    except Exception as e:
        print(f"Error generating response from Gemini: {e}")
        return ["Sorry, I encountered an error while generating a response."] # Return as a list

# --- Chat History Functions ---
def get_chat_history_dir():
    return "chat_history"

def save_chat_history(session_id, messages, doc_info):
    """
    Saves chat messages and associated document information for a session.
    doc_info should be a dict like {"processed_files": [], "vector_store_path": ""}
    """
    history_dir = get_chat_history_dir()
    if not os.path.exists(history_dir):
        os.makedirs(history_dir)
    filepath = os.path.join(history_dir, f"{session_id}.json")
    
    session_data = {
        "messages": messages,
        "doc_info": doc_info
    }

    try:
        with open(filepath, "w") as f:
            json.dump(session_data, f, indent=2)
        # Store timestamp of last message to sort chats
        with open(filepath, "a") as f: # Append a line with timestamp for easy retrieval
            f.write(f"\n# LastModified: {datetime.datetime.now().isoformat()}")
        print(f"Chat history and doc info saved to: {filepath}")
    except Exception as e:
        st.error(f"Error saving chat history: {e}")

def load_chat_history(session_id):
    """
    Loads chat messages and associated document information for a session.
    Handles both old (list of messages) and new (dict with messages and doc_info) formats.
    Returns (messages, doc_info) or ([], default_doc_info) if not found/corrupted.
    """
    history_dir = get_chat_history_dir()
    filepath = os.path.join(history_dir, f"{session_id}.json")
    
    default_doc_info = {"processed_files": [], "vector_store_path": None}

    if not os.path.exists(filepath):
        print(f"Chat history file not found: {filepath}")
        return [], default_doc_info
    
    try:
        with open(filepath, "r") as f:
            lines = f.readlines()
            json_content = "".join([line for line in lines if not line.strip().startswith("# LastModified:")])
            
            loaded_data = json.loads(json_content)
            
            messages = []
            doc_info = default_doc_info

            # Check if the loaded data is an old format (list of messages)
            if isinstance(loaded_data, list):
                messages = loaded_data
                print(f"Loaded old format chat history for session '{session_id}'.")
            # Check if the loaded data is the new format (dictionary)
            elif isinstance(loaded_data, dict):
                messages = loaded_data.get("messages", [])
                doc_info = loaded_data.get("doc_info", default_doc_info)
                print(f"Loaded new format chat history for session '{session_id}'.")
            else:
                raise ValueError("Unexpected chat history file format.")
            
            print(f"Chat history loaded from: {filepath}")
            return messages, doc_info
    except json.JSONDecodeError as e:
        st.error(f"Error loading chat history for session '{session_id}': The file is corrupted or empty. Starting a new history for this session. Details: {e}")
        print(f"JSONDecodeError for session {session_id}: {e}")
        return [], default_doc_info
    except Exception as e:
        st.error(f"An unexpected error occurred while loading chat history for session '{session_id}'. Details: {e}")
        print(f"Error loading chat history for session {session_id}: {e}")
        return [], default_doc_info

def clear_chat_history_file(session_id):
    history_dir = get_chat_history_dir()
    filepath = os.path.join(history_dir, f"{session_id}.json")
    
    success = False
    if os.path.exists(filepath):
        try:
            os.remove(filepath)
            print(f"Chat history file deleted: {filepath}")
            success = True
        except Exception as e:
            st.error(f"Error clearing chat history file: {e}")
    else:
        print(f"Chat history file not found, cannot delete: {filepath}")

    # Also remove the associated FAISS index directory if it exists
    index_path = os.path.join(FAISS_INDICES_DIR, f"{session_id}_index")
    if os.path.exists(index_path):
        try:
            shutil.rmtree(index_path) 
            print(f"FAISS index directory deleted: {index_path}")
        except Exception as e:
            st.error(f"Error deleting FAISS index: {e}")
            
    return success

def get_or_create_session_id():
    if "session_id" not in st.session_state:
        st.session_state.session_id = str(uuid.uuid4())
        print(f"Generated new session ID: {st.session_state.session_id}")
    else:
        print(f"Using existing session ID: {st.session_state.session_id}")
    return st.session_state.session_id

# --- get_all_chat_sessions function ---
def get_all_chat_sessions():
    sessions = []
    history_dir = get_chat_history_dir()
    if not os.path.exists(history_dir):
        return sessions

    for filename in os.listdir(history_dir):
        if filename.endswith(".json"):
            session_id = filename.replace(".json", "")
            filepath = os.path.join(history_dir, filename)
            
            last_modified_time = None
            first_user_message_content = ""

            try:
                # Read content to find timestamp and first user message
                with open(filepath, "r") as f:
                    lines = f.readlines()
                    for line in lines:
                        if line.startswith("# LastModified:"):
                            try:
                                last_modified_time = datetime.datetime.fromisoformat(line.split(": ")[1].strip())
                            except ValueError:
                                pass # Handle potential old/malformed timestamps

                    # Extract JSON part
                    json_content = "".join([line for line in lines if not line.strip().startswith("# LastModified:")])
                    
                    loaded_data = json.loads(json_content)
                    
                    history = []
                    # Check if the loaded data is an old format (list of messages)
                    if isinstance(loaded_data, list):
                        history = loaded_data
                    # Check if the loaded data is the new format (dictionary)
                    elif isinstance(loaded_data, dict):
                        history = loaded_data.get("messages", [])
                    else:
                        # If it's neither, skip this file
                        print(f"Skipping file {filename} due to unexpected format.")
                        continue

                    for message in history:
                        if message["role"] == "user" and message["content"]:
                            first_user_message_content = message["content"].split('\n')[0][:50]
                            if len(message["content"].split('\n')[0]) > 50:
                                first_user_message_content += "..."
                            break

            except (json.JSONDecodeError, FileNotFoundError, Exception) as e:
                print(f"Error processing chat history file {filename}: {e}")
                continue # Skip corrupted files

            if not last_modified_time:
                # If no timestamp found in file, use file's last modified time as a fallback
                last_modified_time = datetime.datetime.fromtimestamp(os.path.getmtime(filepath))

            sessions.append({
                "id": session_id,
                "title": first_user_message_content if first_user_message_content else f"Chat {session_id[:8]}",
                "last_modified": last_modified_time
            })
            
    # Sort sessions by last modified time, newest first
    sessions.sort(key=lambda x: x["last_modified"], reverse=True)
    return sessions

# --- Helper for categorizing sessions ---
def get_category_for_date(chat_date):
    today = datetime.date.today()
    if chat_date == today:
        return "Today"
    elif chat_date == today - datetime.timedelta(days=1):
        return "Yesterday"
    elif today - datetime.timedelta(days=7) <= chat_date < today - datetime.timedelta(days=1):
        return "Previous 7 Days"
    elif today - datetime.timedelta(days=30) <= chat_date < today - datetime.timedelta(days=7):
        return "Previous 30 Days"
    else:
        return chat_date.strftime("%B %Y") # Month Year for Older chats

# --- Process User Question (Modified to use session-specific index) ---
def process_user_question(user_question, selected_model):
    context_text = ""
    vector_store_path = st.session_state.current_doc_info.get("vector_store_path")
    
    # Check if there's a vector store path for the current session
    if vector_store_path and os.path.exists(vector_store_path):
        try:
            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            new_db = FAISS.load_local(
                vector_store_path, embeddings, allow_dangerous_deserialization=True
            )
            docs = new_db.similarity_search(user_question, k=3)
            context_text = "\n\n".join([doc.page_content for doc in docs])
            
            print(f"RAG Attempt: Context retrieved (length {len(context_text)}).")
            st.session_state.vector_store_exists = True # Flag that a vector store is active for RAG
        except ValueError as e:
            st.warning(f"FAISS index at '{vector_store_path}' not found or corrupted, proceeding with general knowledge. Details: {e}")
            st.session_state.vector_store_exists = False
        except Exception as e:
            st.warning(f"Error during document retrieval from '{vector_store_path}', proceeding with general knowledge. Details: {e}")
            st.session_state.vector_store_exists = False
    else:
        print("No documents processed for this session or vector store not found. Using general chat only.")
        st.session_state.vector_store_exists = False # No vector store active for RAG

    raw_llm_response_generator = get_llm_response(
        user_question, context_text, selected_model, st.session_state.vector_store_exists
    )
    return raw_llm_response_generator

# --- Main Streamlit Application ---
def main():
    st.set_page_config("Chat with LLMs")
    st.title("Analyse Your Documents & Chat with LLMs 💁")

    # Initialize session states for new/first load
    if "session_id" not in st.session_state:
        st.session_state.session_id = str(uuid.uuid4())
        st.session_state.chat_history = []
        st.session_state.current_doc_info = {"processed_files": [], "vector_store_path": None}
        st.session_state.vector_store_exists = False
        st.session_state.loaded_session_id = None # To track if current session's data has been loaded
        st.session_state.accumulated_uploaded_files_data = {} # Stores file content for processing (new)

    # Logic to load chat history and doc info when session_id changes or on first load
    if st.session_state.session_id != st.session_state.get("loaded_session_id"):
        current_session_id = st.session_state.session_id
        messages, doc_info = load_chat_history(current_session_id)
        
        st.session_state.chat_history = messages
        st.session_state.current_doc_info = doc_info
        st.session_state.vector_store_exists = True if doc_info.get("vector_store_path") and os.path.exists(doc_info["vector_store_path"]) else False
        st.session_state.loaded_session_id = current_session_id # Mark this session as loaded
        
        # When loading a session, clear the accumulated_uploaded_files_data
        # as the actual file objects cannot be easily persisted. User will need to
        # re-upload files if they want to modify the document set for a loaded session.
        st.session_state.accumulated_uploaded_files_data = {} 
        
        print(f"Initialized/Loaded chat history for session {current_session_id}: {len(st.session_state.chat_history)} messages, Doc Info: {st.session_state.current_doc_info}")
    else:
        print(f"Using existing chat history for session {st.session_state.session_id}: {len(st.session_state.chat_history)} messages, Doc Info: {st.session_state.current_doc_info}")


    if "widget_states" not in st.session_state:
        st.session_state.widget_states = {}

    # --- Sidebar ---
    with st.sidebar:
        st.title("Settings")

        # Model selection dropdown
        model_options = ["Llama 3.2 (Ollama)", "Gemini API"]
        selected_model = st.selectbox("Choose your LLM:", model_options, key="model_selector")

        # "New Chat" button
        if st.button("➕ New Chat", key="new_chat_button"):
            # Clear all session data and create a new session ID
            st.session_state.session_id = str(uuid.uuid4())
            st.session_state.chat_history = []
            st.session_state.current_doc_info = {"processed_files": [], "vector_store_path": None}
            st.session_state.vector_store_exists = False
            st.session_state.widget_states["text_input"] = "" # Clear main chat input
            st.session_state.loaded_session_id = None # Force reload on next rerun
            st.session_state.accumulated_uploaded_files_data = {} # Clear accumulated files for new chat
            st.rerun()

        st.markdown("---")
        st.subheader("Recent Chats")
        
        all_sessions = get_all_chat_sessions()
        
        if "show_all_chats" not in st.session_state:
            st.session_state.show_all_chats = False

        DISPLAY_LIMIT = 4 # Number of recent chats to show initially

        chats_to_display = []
        if st.session_state.show_all_chats:
            chats_to_display = all_sessions
        else:
            chats_to_display = all_sessions[:DISPLAY_LIMIT]

        last_displayed_category = None
        
        if chats_to_display:
            for session in chats_to_display:
                current_category = get_category_for_date(session["last_modified"].date())
                
                if current_category != last_displayed_category:
                    if current_category in ["Today", "Yesterday", "Previous 7 Days", "Previous 30 Days"]:
                        st.markdown(f"**{current_category}**")
                    else:
                        st.markdown(f"***{current_category}***")
                    last_displayed_category = current_category
                
                if st.button(session["title"], key=f"load_chat_{session['id']}"):
                    st.session_state.session_id = session["id"] # Set current session_id
                    st.session_state.loaded_session_id = None # Force reload on next rerun to load specific chat data
                    st.session_state.widget_states["text_input"] = ""
                    st.session_state.accumulated_uploaded_files_data = {} # Clear uploader cache on chat switch
                    st.rerun()

            # The "Show more/less" button logic
            if len(all_sessions) > DISPLAY_LIMIT:
                if st.session_state.show_all_chats:
                    if st.button(f"Show less 🔼", key=f"show_less_all_chats"):
                        st.session_state.show_all_chats = False
                        st.rerun()
                else:
                    if st.button(f"Show more 🔽", key=f"show_more_all_chats"):
                        st.session_state.show_all_chats = True
                        st.rerun()
            st.markdown("---")
        else:
            st.info("No recent chats found.")
            st.markdown("---")

        st.subheader("Document Upload & Processing")
        
        uploaded_files_this_run = st.file_uploader(
            "Upload more documents to add to current session:",
            accept_multiple_files=True,
            type=["pdf", "pptx", "ppt", "xlsx", "xls"],
            key=f"document_uploader_{st.session_state.session_id}" # Stable key for session
        )
        
        # Add newly uploaded files to the accumulated list
        if uploaded_files_this_run:
            for uploaded_file in uploaded_files_this_run:
                # Store the file content by name and size to handle unique names/multiple uploads
                file_unique_id = (uploaded_file.name, uploaded_file.size)
                if file_unique_id not in st.session_state.accumulated_uploaded_files_data:
                    st.session_state.accumulated_uploaded_files_data[file_unique_id] = uploaded_file

        # Process Documents button
        if st.button("Process ALL Uploaded Documents", key="process_docs_button"):
            # Get the list of files to process from the accumulated data
            files_to_process = list(st.session_state.accumulated_uploaded_files_data.values())

            if files_to_process:
                with st.spinner("Extracting text and creating document index..."):
                    all_text_list = [] # Use a list to accumulate text from each successfully processed file
                    successfully_processed_file_names = [] # Track only successfully processed files
                    
                    for uploaded_file in files_to_process:
                        file_type = uploaded_file.name.split(".")[-1].lower()
                        file_specific_text = "" # Text extracted from the current file
                        try:
                            # IMPORTANT: Reset stream position before reading each file.
                            # This is crucial if the file object has been read before
                            # (e.g., in a previous 'st.file_uploader' interaction)
                            uploaded_file.seek(0) 

                            if file_type == "pdf":
                                file_specific_text = get_pdf_text([uploaded_file])
                            elif file_type in ["pptx", "ppt"]:
                                prs = Presentation(uploaded_file)
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if shape.has_text_frame:
                                            text_frame = shape.text_frame
                                            for paragraph in text_frame.paragraphs:
                                                for run in paragraph.runs:
                                                    file_specific_text += run.text + " "
                            elif file_type in ["xlsx", "xls"]:
                                excel_data = pd.read_excel(uploaded_file, sheet_name=None)
                                for sheet_name, df in excel_data.items():
                                    file_specific_text += df.to_string(index=False, header=True) + "\n\n"
                            else:
                                st.warning(f"File type {file_type} of '{uploaded_file.name}' is not supported and will be skipped.")
                                continue # Skip to next file

                            if file_specific_text.strip(): # Check if any text was actually extracted
                                all_text_list.append(file_specific_text)
                                successfully_processed_file_names.append(uploaded_file.name)
                            else:
                                st.warning(f"Could not extract any text from '{uploaded_file.name}'. It will not be included in the context.")

                        except Exception as e:
                            st.error(f"Error processing '{uploaded_file.name}': {e}. This file will be skipped.")
                            continue # Continue to the next file even if one fails

                    # Combine all successfully extracted text into one string
                    all_combined_text = "\n\n".join(all_text_list)

                    if all_combined_text.strip() and successfully_processed_file_names:
                        text_chunks = get_text_chunks(all_combined_text)
                        vector_store_path = get_vector_store(text_chunks, st.session_state.session_id)
                        
                        st.session_state.current_doc_info = {
                            "processed_files": successfully_processed_file_names,
                            "vector_store_path": vector_store_path
                        }
                        st.session_state.vector_store_exists = True
                        st.success(f"Documents processed and ready for intelligent Q&A! ({len(successfully_processed_file_names)} files included in context)")
                    else:
                        st.warning("No extractable text found in any of the uploaded documents. General chat only.")
                        st.session_state.current_doc_info = {"processed_files": [], "vector_store_path": None}
                        st.session_state.vector_store_exists = False
            else:
                st.warning("No documents selected to process. Please upload files first.")
                st.session_state.current_doc_info = {"processed_files": [], "vector_store_path": None}
                st.session_state.vector_store_exists = False

        # Display currently associated documents (those whose content contributed to the vector store)
        if st.session_state.current_doc_info.get("processed_files"):
            st.markdown("---")
            st.subheader("Currently Processed Documents:")
            for doc_name in st.session_state.current_doc_info["processed_files"]:
                st.markdown(f"- {doc_name}")
            st.info(f"Vector store active: {st.session_state.vector_store_exists}")
        else:
            st.info(f"No documents processed for this session. Vector store active: {st.session_state.vector_store_exists}")

        # Display list of documents that have been uploaded but not necessarily processed yet
        if st.session_state.accumulated_uploaded_files_data:
            st.markdown("---")
            st.subheader("Files Uploaded for Current Session (Pending Processing):")
            for file_unique_id in st.session_state.accumulated_uploaded_files_data.keys():
                st.markdown(f"- {file_unique_id[0]}") # Display only the filename
            if st.button("Clear All Uploaded Documents", key="clear_all_uploaded_docs_button"):
                st.session_state.accumulated_uploaded_files_data = {}
                st.session_state.current_doc_info = {"processed_files": [], "vector_store_path": None}
                st.session_state.vector_store_exists = False
                # Remove FAISS index for current session if it exists
                index_path = os.path.join(FAISS_INDICES_DIR, f"{st.session_state.session_id}_index")
                if os.path.exists(index_path):
                    try:
                        shutil.rmtree(index_path) 
                        print(f"FAISS index directory deleted: {index_path}")
                    except Exception as e:
                        st.error(f"Error deleting FAISS index: {e}")
                st.success("All uploaded documents and associated vector store cleared for this session.")
                st.rerun() # Rerun to update the display

        st.markdown("---")
        if st.button("Clear Current Chat & Documents", key="clear_current_chat_button"):
            if clear_chat_history_file(st.session_state.session_id):
                # Reset all session states after clearing the file
                st.session_state.session_id = str(uuid.uuid4()) # Generate new ID
                st.session_state.chat_history = []
                st.session_state.current_doc_info = {"processed_files": [], "vector_store_path": None}
                st.session_state.vector_store_exists = False
                st.session_state.widget_states["text_input"] = ""
                st.session_state.loaded_session_id = None # Force reload on next rerun
                st.session_state.accumulated_uploaded_files_data = {} # Clear accumulated files
                st.rerun()
                st.success("Current chat history and associated documents cleared!")
            else:
                st.error("Failed to clear current chat history.")

    # Display chat history
    if st.session_state.chat_history:
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

    # Main chat input and response
    prompt = st.chat_input("Ask a question...", key="text_input")
    if prompt:
        st.session_state.chat_history.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            
            response_generator = process_user_question(prompt, selected_model)
            
            full_response = ""
            for chunk in response_generator:
                full_response += chunk
                message_placeholder.markdown(full_response + " ▌")

            message_placeholder.markdown(full_response)
            st.session_state.chat_history.append({"role": "assistant", "content": full_response})

        # Save history AFTER the response is generated, including doc_info
        save_chat_history(
            st.session_state.session_id,
            st.session_state.chat_history,
            st.session_state.current_doc_info
        )
        # Rerun to update the sidebar with the new chat history and its ordering
        st.rerun()
    
if __name__ == "__main__":
    try:
        # Ensure chat_history and faiss_indices directories exist
        if not os.path.exists(get_chat_history_dir()):
            os.makedirs(get_chat_history_dir())
        if not os.path.exists(FAISS_INDICES_DIR):
            os.makedirs(FAISS_INDICES_DIR)
            
        main()
    except Exception as e:
        st.error(f"An unexpected error occurred: {e}")
        print(f"Error executing main: {e}")